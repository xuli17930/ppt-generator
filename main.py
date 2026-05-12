from fastapi import FastAPI
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, uuid, requests
from io import BytesIO

app = FastAPI()

class SlideContent(BaseModel):
    title: str
    bullets: List[str]

class PPTRequest(BaseModel):
    main_title: str
    subtitle: str
    slides: List[SlideContent]
    cover_image: str = ""
    data_images: List[str] = []
    data_points: List[Dict] = []
    summary: str = ""

# 配色方案
DARK_BLUE = RGBColor(0, 32, 96)
BRIGHT_BLUE = RGBColor(0, 112, 192)
LIGHT_BLUE = RGBColor(230, 240, 255)
TEXT_DARK = RGBColor(33, 33, 33)
TEXT_GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)

def download_image(url):
    try:
        response = requests.get(url, timeout=15, headers={'User-Agent': 'Mozilla/5.0'})
        if response.status_code == 200:
            return BytesIO(response.content)
    except Exception as e:
        print(f"图片下载失败: {e}")
    return None

@app.post("/generate")
def generate_ppt(req: PPTRequest):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

       # ========== 封面页（标题与图片不重叠） ==========
    slide = prs.slides.add_slide(blank)
    
    if req.cover_image:
        img_stream = download_image(req.cover_image)
        if img_stream:
            # 图片从 x=6.0 开始，给文字留出 5.2 英寸安全区
            slide.shapes.add_picture(img_stream, Inches(6.0), Inches(0), height=Inches(7.5))
            # 左侧遮罩（淡蓝背景）
            mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6.0), Inches(7.5))
            mask.fill.solid()
            mask.fill.fore_color.rgb = LIGHT_BLUE
            mask.line.fill.background()
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = DARK_BLUE; top_bar.line.fill.background()
    
    # 主标题（宽度限制在 5.0 英寸，自动换行）
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.0), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = req.main_title
    p.font.size = Pt(28)        # 稍微缩小，长标题也能显示
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Microsoft YaHei"
    
    # 副标题
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(5.0), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = req.subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_GRAY
    p2.font.name = "Microsoft YaHei"
    
    # 底部装饰线
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.5), Inches(2.5), Inches(0.03))
    bottom_line.fill.solid(); bottom_line.fill.fore_color.rgb = BRIGHT_BLUE; bottom_line.line.fill.background()

       # ========== 文章综述页（大容量 + 智能分页） ==========
    if req.summary and len(req.summary) > 0:
        paragraphs = [p.strip() for p in req.summary.split('\n') if p.strip()]
        
        # 一页容量约 1200 字（11pt 字体 + 紧凑行距）
        max_chars_per_page = 1200
        pages = []
        current_page = []
        current_count = 0
        
        for para in paragraphs:
            para_len = len(para)
            if current_count + para_len > max_chars_per_page and len(current_page) > 0:
                pages.append(current_page)
                current_page = [para]
                current_count = para_len
            else:
                current_page.append(para)
                current_count += para_len
        
        if len(current_page) > 0:
            pages.append(current_page)
        
        for page_idx, page_paras in enumerate(pages):
            slide = prs.slides.add_slide(blank)
            
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
            
            header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
            header_bg.fill.solid(); header_bg.fill.fore_color.rgb = LIGHT_BLUE; header_bg.line.fill.background()
            
            left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(1.0))
            left_bar.fill.solid(); left_bar.fill.fore_color.rgb = BRIGHT_BLUE; left_bar.line.fill.background()
            
            # 标题
            if len(pages) == 1:
                title_text = "文章综述"
            else:
                title_text = "文章综述（上）" if page_idx == 0 else "文章综述（下）"
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = title_text
            tp.font.size = Pt(24); tp.font.bold = True
            tp.font.color.rgb = DARK_BLUE; tp.font.name = "Microsoft YaHei"
            
            content_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.0))
            content_bg.fill.solid(); content_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            content_bg.line.color.rgb = RGBColor(220, 220, 220)
            
            # 扩大文本框高度到 6.2 英寸
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(6.2))
            tf = text_box.text_frame; tf.word_wrap = True
            
            for i, para in enumerate(page_paras):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para
                p.font.size = Pt(11)        # 11pt 更小字号
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(6)       # 更紧凑行距
                p.font.name = "Microsoft YaHei"
                p.level = 0
            
            footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.35), Inches(12.3), Inches(0.02))
            footer_line.fill.solid(); footer_line.fill.fore_color.rgb = BRIGHT_BLUE; footer_line.line.fill.background()
            
    # ========== 核心数据看板页 ==========
    if req.data_points and len(req.data_points) > 0:
        slide = prs.slides.add_slide(blank)
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BLUE
        bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "核心数据"
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        tp.font.name = "Microsoft YaHei"
        
        card_width = 4.0
        start_x = 0.8
        start_y = 1.5
        gap = 0.3
        
        for i, dp in enumerate(req.data_points[:3]):
            x = start_x + (i % 3) * (card_width + gap)
            y = start_y + (i // 3) * 2.5
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), 
                Inches(card_width), Inches(2.0)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = BRIGHT_BLUE
            
            num_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.2), Inches(card_width - 0.4), Inches(1.0)
            )
            np = num_box.text_frame.paragraphs[0]
            
            if dp.get('type') == 'compare':
                np.text = f"{dp.get('before', '')} → {dp.get('after', '')}"
            else:
                np.text = str(dp.get('value', ''))
            
            np.font.size = Pt(32)
            np.font.bold = True
            np.font.color.rgb = BRIGHT_BLUE
            np.alignment = PP_ALIGN.CENTER
            np.font.name = "Microsoft YaHei"
            
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 1.2), Inches(card_width - 0.4), Inches(0.5)
            )
            lp = label_box.text_frame.paragraphs[0]
            lp.text = dp.get('label', '')
            lp.font.size = Pt(14)
            lp.font.color.rgb = TEXT_GRAY
            lp.alignment = PP_ALIGN.CENTER
            lp.font.name = "Microsoft YaHei"

    # ========== 内容页 ==========
    data_img_idx = 0
    
    for idx, s in enumerate(req.slides, 1):
        slide = prs.slides.add_slide(blank)
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()
        
        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = LIGHT_BLUE
        header_bg.line.fill.background()
        
        left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(1.1))
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = BRIGHT_BLUE
        left_bar.line.fill.background()
        
        page_tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.8), Inches(0.25), Inches(1.0), Inches(0.5))
        page_tag.fill.solid()
        page_tag.fill.fore_color.rgb = BRIGHT_BLUE
        page_tag.line.fill.background()
        
        tag_text = slide.shapes.add_textbox(Inches(11.8), Inches(0.25), Inches(1.0), Inches(0.5))
        ttp = tag_text.text_frame.paragraphs[0]
        ttp.text = f"0{idx}" if idx < 10 else str(idx)
        ttp.font.size = Pt(14)
        ttp.font.bold = True
        ttp.font.color.rgb = WHITE
        ttp.alignment = PP_ALIGN.CENTER
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11.0), Inches(0.7))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = s.title
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = DARK_BLUE
        tp.font.name = "Microsoft YaHei"
        
        title_lower = s.title.lower()
        need_image = any(k in title_lower for k in ['数据', '图表', '实验', '结果', '对比', '分析', '统计', '调研'])
        has_image = need_image and data_img_idx < len(req.data_images)
        
        if has_image:
            content_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(6.8), Inches(5.8))
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            content_bg.line.color.rgb = RGBColor(220, 220, 220)
            
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.4), Inches(5.4))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, b in enumerate(s.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"●  {b}"
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(10)
                p.font.name = "Microsoft YaHei"
            
            img_stream = download_image(req.data_images[data_img_idx])
            if img_stream:
                slide.shapes.add_picture(img_stream, Inches(7.6), Inches(1.6), width=Inches(5.0))
                data_img_idx += 1
        else:
            content_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8))
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
            content_bg.line.color.rgb = RGBColor(220, 220, 220)
            
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(5.4))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, b in enumerate(s.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"●  {b}"
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(12)
                p.font.name = "Microsoft YaHei"
        
        footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.35), Inches(12.3), Inches(0.02))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = BRIGHT_BLUE
        footer_line.line.fill.background()

    # ========== 结尾页 ==========
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    thanks_box = slide.shapes.add_textbox(Inches(0), Inches(2.8), Inches(13.333), Inches(1.2))
    tp = thanks_box.text_frame.paragraphs[0]
    tp.text = "感谢观看"
    tp.font.size = Pt(48)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    tp.font.name = "Microsoft YaHei"

    # 保存
    os.makedirs("/app/output", exist_ok=True)
    # 用文章标题作为文件名，清理非法字符
import re
safe_title = re.sub(r'[\\/:*?"<>|]', '', req.main_title)  # 去掉Windows/Unix非法字符
safe_title = safe_title[:50]  # 限制长度50字，避免文件名过长
fname = f"{safe_title}.pptx"
    fpath = f"/app/output/{fname}"
    prs.save(fpath)
    
    base = os.environ.get("RENDER_EXTERNAL_URL", "")
    if not base:
        base = "https://" + os.environ.get("RENDER_EXTERNAL_HOSTNAME", "")
    return {"download_url": f"{base}/download/{fname}", "filename": fname}

@app.get("/download/{filename}")
def download(filename: str):
    fpath = f"/app/output/{filename}"
    if os.path.exists(fpath):
        return FileResponse(fpath, filename=filename)
    return {"error": "file not found"}
